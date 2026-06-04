import ast
import html
import math
import operator
import os
import re
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any
from urllib.parse import parse_qs, quote_plus, urlparse, unquote

import httpx
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas


GENERATED_DIR = Path(__file__).resolve().parents[1] / "generated"
USER_AGENT = os.getenv("WEB_SEARCH_USER_AGENT", "Mozilla/5.0 (compatible; EthioClaw/1.0)")
WEB_SEARCH_TIMEOUT = int(os.getenv("WEB_SEARCH_TIMEOUT", "20"))
WEB_SEARCH_LIMIT = int(os.getenv("WEB_SEARCH_LIMIT", "5"))


@dataclass
class ToolBundle:
    context: str = ""
    artifact_kind: str | None = None
    artifact_title: str | None = None


def _slugify(value: str) -> str:
    cleaned = re.sub(r"[^a-zA-Z0-9]+", "-", value.lower()).strip("-")
    return cleaned or "artifact"


def _ensure_generated_dir() -> Path:
    GENERATED_DIR.mkdir(parents=True, exist_ok=True)
    return GENERATED_DIR


def _strip_html(value: str) -> str:
    text = re.sub(r"<[^>]+>", "", value or "")
    text = html.unescape(text)
    return re.sub(r"\s+", " ", text).strip()


def sanitize_untrusted_text(value: str, limit: int = 1200) -> str:
    """Remove common prompt-injection phrases from untrusted tool output."""
    text = (value or "").strip()
    if not text:
        return ""

    blocked_patterns = [
        r"ignore previous",
        r"disregard",
        r"system:",
        r"you are now",
        r"forget your instructions",
        r"new instructions",
    ]
    for pattern in blocked_patterns:
        text = re.sub(pattern, "", text, flags=re.I)

    text = re.sub(r"\s+", " ", text).strip()
    return text[:limit]


ALLOWED_BINOPS: dict[type, Any] = {
    ast.Add: operator.add,
    ast.Sub: operator.sub,
    ast.Mult: operator.mul,
    ast.Div: operator.truediv,
    ast.FloorDiv: operator.floordiv,
    ast.Mod: operator.mod,
    ast.Pow: operator.pow,
}

ALLOWED_UNARY: dict[type, Any] = {
    ast.UAdd: operator.pos,
    ast.USub: operator.neg,
}

ALLOWED_FUNCTIONS: dict[str, Any] = {
    "abs": abs,
    "round": round,
    "sqrt": math.sqrt,
    "sin": math.sin,
    "cos": math.cos,
    "tan": math.tan,
    "log": math.log,
    "log10": math.log10,
    "floor": math.floor,
    "ceil": math.ceil,
    "pow": pow,
}

ALLOWED_CONSTANTS: dict[str, float] = {
    "pi": math.pi,
    "e": math.e,
}


def safe_math_eval(expression: str) -> float:
    expression = (expression or "").strip().replace("^", "**")
    tree = ast.parse(expression, mode="eval")

    def _eval(node: ast.AST):
        if isinstance(node, ast.Expression):
            return _eval(node.body)
        if isinstance(node, ast.Constant) and isinstance(node.value, (int, float)):
            return node.value
        if isinstance(node, ast.BinOp) and type(node.op) in ALLOWED_BINOPS:
            return ALLOWED_BINOPS[type(node.op)](_eval(node.left), _eval(node.right))
        if isinstance(node, ast.UnaryOp) and type(node.op) in ALLOWED_UNARY:
            return ALLOWED_UNARY[type(node.op)](_eval(node.operand))
        if isinstance(node, ast.Name) and node.id in ALLOWED_CONSTANTS:
            return ALLOWED_CONSTANTS[node.id]
        if isinstance(node, ast.Call) and isinstance(node.func, ast.Name):
            func = ALLOWED_FUNCTIONS.get(node.func.id)
            if func is None:
                raise ValueError(f"Function '{node.func.id}' is not allowed")
            args = [_eval(arg) for arg in node.args]
            return func(*args)
        raise ValueError("Unsupported math expression")

    return float(_eval(tree))


def _extract_math_expression(message: str) -> str | None:
    text = message.strip()
    lowered = text.lower()

    if not any(
        keyword in lowered
        for keyword in [
            "calculate",
            "compute",
            "what is",
            "solve",
            "how much is",
            "average",
            "mean",
            "sum of",
            "square root",
            "sqrt",
            "percent",
        ]
    ) and not re.search(r"[0-9]\s*[\+\-\*\/\^]\s*[0-9]", text):
        return None

    candidate = text
    for prefix in [
        "calculate",
        "compute",
        "what is",
        "solve",
        "how much is",
    ]:
        if lowered.startswith(prefix):
            candidate = text[len(prefix):].strip(" ?:=.-")
            break

    if "average of" in lowered:
        numbers = re.findall(r"-?\d+(?:\.\d+)?", text)
        if numbers:
            expr = f"({'+'.join(numbers)})/{len(numbers)}"
            return expr

    if "sum of" in lowered:
        numbers = re.findall(r"-?\d+(?:\.\d+)?", text)
        if numbers:
            return "+".join(numbers)

    match = re.search(r"([-+*/().\d\s^eEpiabsqrtlogncof]+)$", candidate)
    if match:
        expr = match.group(1).strip()
        if expr:
            return expr

    math_chars = re.findall(r"[0-9\.\+\-\*\/\^\(\)\s]+", candidate)
    if math_chars:
        expr = "".join(math_chars).strip()
        if expr:
            return expr

    return None


def run_math_tool(message: str) -> str | None:
    expression = _extract_math_expression(message)
    if not expression:
        return None

    try:
        value = safe_math_eval(expression)
        if float(value).is_integer():
            value_text = str(int(value))
        else:
            value_text = str(round(value, 8))
        return f"Math tool: {expression} = {value_text}"
    except Exception as exc:
        return f"Math tool error for '{expression}': {exc}"


def _extract_search_query(message: str) -> str | None:
    text = message.strip()
    lowered = text.lower()
    search_triggers = [
        "search",
        "look up",
        "look for",
        "find",
        "latest",
        "news",
        "web",
        "browse",
        "research",
    ]
    if not any(trigger in lowered for trigger in search_triggers):
        return None

    query = text
    for prefix in [
        "search for",
        "look up",
        "look for",
        "find",
        "research",
        "browse",
    ]:
        if lowered.startswith(prefix):
            query = text[len(prefix):].strip(" :.-")
            break

    query = re.sub(r"^(latest|news on|web search on|search the web for)\s+", "", query, flags=re.I).strip()
    return query or text


def _looks_like_memory_question(message: str) -> bool:
    text = message.strip().lower()
    if not text:
        return False

    memory_phrases = (
        "what did i just ask",
        "what did i ask you",
        "what did i say",
        "what was my last question",
        "what was the last thing i said",
        "what were we talking about",
        "what were we doing",
        "do you remember",
        "remember our conversation",
        "recall our conversation",
        "earlier in this chat",
        "previous message",
        "last message",
    )

    return any(phrase in text for phrase in memory_phrases)


def _should_search_web(message: str) -> bool:
    if _looks_like_memory_question(message):
        return False

    text = message.strip().lower()
    if not text:
        return False

    explicit_search_triggers = (
        "search",
        "look up",
        "look for",
        "find",
        "research",
        "browse",
        "web",
    )

    time_sensitive_triggers = (
        "latest",
        "news",
        "current",
        "today",
        "this week",
        "this month",
        "recent",
        "updated",
        "now",
        "trend",
    )

    if any(trigger in text for trigger in explicit_search_triggers):
        return True

    return any(trigger in text for trigger in time_sensitive_triggers)


def resolve_memory_question(message: str, recent_turns: list[dict[str, str]] | None = None, working_summary: str = "") -> str | None:
    if not _looks_like_memory_question(message):
        return None

    turns = recent_turns or []
    prior_user_turns = [turn for turn in turns if turn.get("role") == "user" and turn.get("content")]
    last_user_message = prior_user_turns[-1]["content"].strip() if prior_user_turns else ""

    text = message.strip().lower()

    if "what did i just ask" in text or "what did i ask you" in text or "what did i say" in text or "last message" in text or "previous message" in text:
        if last_user_message:
            return f"You just asked: \"{last_user_message}\""
        return "I don't have a previous user message in this chat yet."

    if "what were we doing" in text or "what were we talking about" in text or "do you remember" in text or "remember our conversation" in text or "recall our conversation" in text:
        summary = (working_summary or "").strip()
        if summary:
            return summary
        if last_user_message:
            return f"We were talking about: {last_user_message}"
        return "I don't have enough recent chat context yet."

    if last_user_message:
        return f"The last thing you said was: \"{last_user_message}\""

    return "I don't have a previous user message in this chat yet."


def resolve_identity_question(message: str, identity_facts: dict[str, str] | None = None) -> str | None:
    text = message.strip().lower()
    if not text:
        return None

    identity_facts = identity_facts or {}

    def _first_value(keys: list[str]) -> str | None:
        for key in keys:
            value = identity_facts.get(key)
            if value:
                return str(value).strip()
        return None

    if any(phrase in text for phrase in ["what is my name", "what's my name", "do you know my name", "who am i"]):
        value = _first_value(["name"])
        return f"Your name is {value}." if value else "I don't have your name stored yet."

    if any(phrase in text for phrase in ["what is my age", "how old am i", "what's my age"]):
        value = _first_value(["age"])
        return f"Your age is {value}." if value else "I don't have your age stored yet."

    if any(phrase in text for phrase in ["what is my profession", "what do i do", "what is my job", "what's my job"]):
        value = _first_value(["profession", "job", "occupation"])
        return f"Your profession is {value}." if value else "I don't have your profession stored yet."

    if any(phrase in text for phrase in ["what is my nationality", "where am i from", "what's my nationality"]):
        value = _first_value(["nationality", "location"])
        return f"Your nationality/location is {value}." if value else "I don't have your nationality stored yet."

    return None


def _extract_main_body(html_text: str) -> str:
    # 1. Get body if present
    body_match = re.search(r"<body[^>]*>([\s\S]*?)<\/body>", html_text, re.I)
    content = body_match.group(1) if body_match else html_text

    # 2. Strip tags that are navigation, footers, scripts, styles, forms, ads, etc.
    content = re.sub(r"<(script|style|noscript|nav|footer|header|aside|form|iframe)[^>]*>([\s\S]*?)<\/\1>", " ", content, flags=re.I)
    
    # 3. Strip any remaining HTML tags
    content = re.sub(r"<[^>]+>", " ", content)
    
    # 4. Unescape HTML entities
    content = html.unescape(content)
    
    # 5. Clean up whitespace
    content = re.sub(r"\s+", " ", content).strip()
    return content


async def search_web(query: str, limit: int = WEB_SEARCH_LIMIT) -> list[dict[str, str]]:
    if not query:
        return []

    url = "https://lite.duckduckgo.com/lite/"
    async with httpx.AsyncClient(timeout=WEB_SEARCH_TIMEOUT, headers={"User-Agent": USER_AGENT}) as client:
        response = await client.get(url, params={"q": query})
    response.raise_for_status()

    html_text = response.text
    results: list[dict[str, str]] = []
    seen: set[str] = set()

    for match in re.finditer(r'<a[^>]+href="(?P<url>[^"]+)"[^>]*>(?P<title>.*?)</a>', html_text, re.I | re.S):
        raw_url = html.unescape(match.group("url"))
        title = _strip_html(match.group("title"))

        if not title:
            continue

        parsed = urlparse(raw_url)
        final_url = raw_url
        if "duckduckgo.com" in parsed.netloc and parsed.path.endswith("/l/"):
            query_params = parse_qs(parsed.query)
            if "uddg" in query_params:
                final_url = unquote(query_params["uddg"][0])

        if not final_url.startswith("http"):
            continue
        if title.lower() in {"next", "more", "images", "videos", "news", "maps", "shopping"}:
            continue
        if final_url in seen:
            continue
        seen.add(final_url)

        results.append({"title": sanitize_untrusted_text(title, 180), "url": sanitize_untrusted_text(final_url, 500)})
        if len(results) >= limit:
            break

    # For the top 2 results, fetch the full page content using httpx.AsyncClient
    for idx, result in enumerate(results):
        if idx >= 2:
            break
        try:
            # Set a per-article fetch timeout of 5 seconds.
            async with httpx.AsyncClient(timeout=5.0, headers={"User-Agent": USER_AGENT}) as client:
                resp = await client.get(result["url"])
                resp.raise_for_status()
                body = _extract_main_body(resp.text)
                if body:
                    result["content"] = body[:1500]
        except Exception as e:
            print(f"⚠️ Failed to fetch content for {result['url']}: {e}")

    return results


def format_search_results(results: list[dict[str, str]], query: str) -> str:
    if not results:
        return sanitize_untrusted_text(f"Web search: no results found for '{query}'.")

    lines = [sanitize_untrusted_text(f"Web search results for '{query}':")]
    for idx, result in enumerate(results, 1):
        content = result.get("content")
        if content:
            lines.append(sanitize_untrusted_text(f"[{idx}] Source: {result['url']}"))
            lines.append(sanitize_untrusted_text(f"Title: {result['title']}"))
            lines.append(sanitize_untrusted_text(f"Content: {content}"))
            lines.append("")
        else:
            lines.append(sanitize_untrusted_text(f"{idx}. {result['title']} - {result['url']}"))
    return "\n".join(lines)


def _extract_artifact_request(message: str) -> tuple[str | None, str | None]:
    lowered = message.lower()
    if any(keyword in lowered for keyword in ["powerpoint", "ppt", "pptx", "presentation", "slides"]):
        return "pptx", _extract_topic(message)
    if any(keyword in lowered for keyword in ["word", "docx", "document", "report"]):
        return "docx", _extract_topic(message)
    if "pdf" in lowered:
        return "pdf", _extract_topic(message)
    return None, None


def _extract_topic(message: str) -> str:
    text = message.strip()
    cleaned = re.sub(
        r"(?i)\b(create|make|generate|draft|build|write|export|prepare|turn this into|turn into|a|an|the|to|into|for)\b",
        " ",
        text,
    )
    cleaned = re.sub(r"(?i)\b(powerpoint|pptx|ppt|presentation|slides|word|docx|document|report|pdf)\b", " ", cleaned)
    cleaned = re.sub(r"\s+", " ", cleaned).strip(" .:-")
    return cleaned or "generated-content"


def clean_markdown(text: str) -> str:
    """Strips bold, italic, and header markers from text."""
    if not text:
        return ""
    # Remove headers (### )
    text = re.sub(r"^#+\s*", "", text, flags=re.MULTILINE)
    # Remove bold/italic (**, __, *, _)
    text = re.sub(r"(\*\*|__)(.*?)\1", r"\2", text)
    text = re.sub(r"(\*|_)(.*?)\1", r"\2", text)
    # Remove list markers from start of line
    text = re.sub(r"^[ \t]*[-*+]\s+", "", text, flags=re.MULTILINE)
    return text.strip()


def _split_paragraphs(content: str) -> list[str]:
    parts = [chunk.strip() for chunk in re.split(r"\n\s*\n", content or "") if chunk.strip()]
    if parts:
        return [clean_markdown(p) for p in parts]
    lines = [line.strip() for line in (content or "").splitlines() if line.strip()]
    if lines:
        return [clean_markdown(l) for l in lines]
    return [clean_markdown(content.strip()) or "No content provided."]


def generate_docx(title: str, content: str) -> Path:
    _ensure_generated_dir()
    file_path = GENERATED_DIR / f"{_slugify(title)}-{datetime.utcnow().strftime('%Y%m%d%H%M%S')}.docx"

    document = Document()
    document.add_heading(clean_markdown(title), level=1)
    for paragraph in _split_paragraphs(content):
        document.add_paragraph(paragraph)
    document.save(file_path)
    return file_path


def generate_pdf(title: str, content: str) -> Path:
    _ensure_generated_dir()
    file_path = GENERATED_DIR / f"{_slugify(title)}-{datetime.utcnow().strftime('%Y%m%d%H%M%S')}.pdf"
    c = canvas.Canvas(str(file_path), pagesize=letter)
    width, height = letter
    y = height - 72

    c.setFont("Helvetica-Bold", 16)
    c.drawString(72, y, clean_markdown(title)[:90])
    y -= 28
    c.setFont("Helvetica", 11)

    for paragraph in _split_paragraphs(content):
        for line in _wrap_text(paragraph, max_chars=95):
            if y < 72:
                c.showPage()
                c.setFont("Helvetica", 11)
                y = height - 72
            c.drawString(72, y, line)
            y -= 14
        y -= 8

    c.save()
    return file_path


def _wrap_text(text: str, max_chars: int = 90) -> list[str]:
    words = text.split()
    lines: list[str] = []
    current = []
    current_len = 0

    for word in words:
        addition = len(word) + (1 if current else 0)
        if current_len + addition > max_chars:
            lines.append(" ".join(current))
            current = [word]
            current_len = len(word)
        else:
            current.append(word)
            current_len += addition

    if current:
        lines.append(" ".join(current))

    return lines or [text]


def generate_pptx(title: str, content: str) -> Path:
    _ensure_generated_dir()
    file_path = GENERATED_DIR / f"{_slugify(title)}-{datetime.utcnow().strftime('%Y%m%d%H%M%S')}.pptx"

    presentation = Presentation()

    # Create a cohesive visual system instead of relying on the default template.
    colors = {
        "navy": "0F172A",
        "ink": "111827",
        "slate": "334155",
        "cyan": "22D3EE",
        "emerald": "34D399",
        "gold": "FBBF24",
        "white": "F8FAFC",
    }

    def apply_background(slide, fill_color: str = colors["navy"]) -> None:
        background = slide.background
        background.fill.solid()
        background.fill.fore_color.rgb = hex_to_rgb(fill_color)

    def add_accent_bar(slide, y: float = 0.28) -> None:
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0),
            Inches(y),
            Inches(13.33),
            Inches(0.08),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(colors["cyan"])
        bar.line.color.rgb = hex_to_rgb(colors["cyan"])
        bar.line.width = Pt(0.5)

    def add_footer(slide, text: str = "Generated by EthioClaw") -> None:
        footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(4.5), Inches(0.25))
        frame = footer.text_frame
        frame.clear()
        p = frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(10)
        p.font.color.rgb = hex_to_rgb(colors["slate"])
        p.font.name = "Aptos"

    def add_title(slide, text: str, top: float = 0.55, size: int = 28) -> None:
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.0), Inches(0.8))
        frame = title_box.text_frame
        frame.clear()
        p = frame.paragraphs[0]
        p.text = text
        p.font.bold = True
        p.font.size = Pt(size)
        p.font.color.rgb = hex_to_rgb(colors["white"])
        p.font.name = "Aptos Display"

    def add_subtitle(slide, text: str, top: float = 1.25) -> None:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(11.8), Inches(0.7))
        frame = box.text_frame
        frame.clear()
        p = frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.color.rgb = hex_to_rgb(colors["cyan"])
        p.font.name = "Aptos"

    def add_card(slide, left: float, top: float, width: float, height: float, fill: str, border: str) -> Any:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill)
        shape.line.color.rgb = hex_to_rgb(border)
        shape.line.width = Pt(1.2)
        return shape

    def add_bullets(slide, bullets: list[str], left: float, top: float, width: float, height: float) -> None:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        frame.clear()
        for idx, bullet in enumerate(bullets):
            p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
            p.text = bullet
            p.level = 0
            p.bullet = True
            p.font.size = Pt(20)
            p.font.color.rgb = hex_to_rgb(colors["white"])
            p.font.name = "Aptos"

    def split_content_slides(text: str) -> list[dict]:
        """Parses content into structured slides by looking for headers or double newlines."""
        # Detect if the LLM used "Slide X:" or "### Slide X:" markers
        slide_pattern = r"(?i)(?:###\s*)?Slide\s+\d+[:\s-]*"
        parts = re.split(slide_pattern, text or "")
        
        slides = []
        # If we found explicit slide markers
        if len(parts) > 1:
            for part in parts:
                content = part.strip()
                if not content:
                    continue
                lines = content.splitlines()
                title_text = lines[0].strip()
                body_content = "\n".join(lines[1:]).strip()
                slides.append({
                    "title": clean_markdown(title_text),
                    "content": body_content
                })
        else:
            # Fallback to splitting by double newline
            sections = [chunk.strip() for chunk in re.split(r"\n\s*\n", text or "") if chunk.strip()]
            for section in sections:
                lines = section.splitlines()
                title_text = lines[0].strip()
                body_content = "\n".join(lines[1:]).strip()
                slides.append({
                    "title": clean_markdown(title_text),
                    "content": body_content
                })
        
        return slides

    def make_bullets(section: str) -> list[str]:
        raw_lines = section.splitlines()
        bullets = []
        for line in raw_lines:
            cleaned = clean_markdown(line)
            if cleaned:
                bullets.append(cleaned)
        
        if len(bullets) > 1:
            return bullets
            
        # Split a long paragraph into natural bullet-sized chunks.
        text = clean_markdown(section)
        sentence_parts = re.split(r"(?<=[.!?])\s+", text)
        return [part.strip() for part in sentence_parts if part.strip()] or [text]

    def hex_to_rgb(hex_value: str):
        from pptx.dml.color import RGBColor

        value = hex_value.lstrip("#")
        return RGBColor.from_string(value)

    # Initial Title Slide
    title_slide = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide)
    apply_background(slide, colors["navy"])
    add_title(slide, clean_markdown(title), top=0.8, size=30)
    add_subtitle(slide, "A professional presentation generated from AI insights", top=1.55)
    add_accent_bar(slide, y=2.05)

    hero = add_card(slide, 0.6, 2.45, 12.1, 3.55, colors["ink"], colors["cyan"])
    hero_text = hero.text_frame
    hero_text.clear()
    intro = hero_text.paragraphs[0]
    intro.text = "Overview"
    intro.font.size = Pt(18)
    intro.font.bold = True
    intro.font.color.rgb = hex_to_rgb(colors["gold"])
    intro.font.name = "Aptos Display"

    summary_lines = make_bullets(content)[:4]
    for line in summary_lines:
        p = hero_text.add_paragraph()
        p.text = line[:180]
        p.level = 0
        p.bullet = True
        p.font.size = Pt(18)
        p.font.color.rgb = hex_to_rgb(colors["white"])
        p.font.name = "Aptos"

    add_footer(slide)

    # Parse and generate content slides
    slides_data = split_content_slides(content)
    
    # Agenda Slide
    if slides_data:
        overview = presentation.slides.add_slide(presentation.slide_layouts[6])
        apply_background(overview, colors["navy"])
        add_title(overview, "Key Topics", top=0.55, size=26)
        add_accent_bar(overview, y=1.1)
        agenda_card = add_card(overview, 0.75, 1.5, 11.9, 5.1, colors["ink"], colors["slate"])
        agenda_frame = agenda_card.text_frame
        agenda_frame.clear()
        p0 = agenda_frame.paragraphs[0]
        p0.text = "Agenda"
        p0.font.size = Pt(18)
        p0.font.bold = True
        p0.font.color.rgb = hex_to_rgb(colors["cyan"])
        p0.font.name = "Aptos Display"
        for s in slides_data[:6]:
            p = agenda_frame.add_paragraph()
            p.text = s["title"][:100]
            p.bullet = True
            p.font.size = Pt(18)
            p.font.color.rgb = hex_to_rgb(colors["white"])
            p.font.name = "Aptos"
        add_footer(overview)

    # Content Slides
    for index, s_data in enumerate(slides_data, start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        apply_background(slide, colors["navy"] if index % 2 else colors["ink"])
        add_title(slide, s_data["title"][:70], top=0.55, size=24)
        add_accent_bar(slide, y=1.1)
        content_card = add_card(slide, 0.75, 1.45, 11.85, 5.3, colors["ink"], colors["cyan"])
        content_frame = content_card.text_frame
        content_frame.clear()
        
        bullets = make_bullets(s_data["content"])[:6]
        for bullet_index, line in enumerate(bullets):
            p = content_frame.paragraphs[0] if bullet_index == 0 else content_frame.add_paragraph()
            p.text = line[:220]
            p.bullet = True
            p.font.size = Pt(18)
            p.font.color.rgb = hex_to_rgb(colors["white"])
            p.font.name = "Aptos"
        add_footer(slide, f"Generated by EthioClaw | Slide {index + 1}")

    presentation.save(file_path)
    return file_path


def generate_artifact(kind: str, title: str, content: str) -> Path:
    kind = kind.lower()
    if kind == "docx":
        return generate_docx(title, content)
    if kind == "pdf":
        return generate_pdf(title, content)
    if kind == "pptx":
        return generate_pptx(title, content)
    raise ValueError(f"Unsupported artifact kind: {kind}")


async def build_tool_bundle(message: str) -> ToolBundle:
    context_parts: list[str] = []

    math_result = run_math_tool(message)
    if math_result:
        context_parts.append(math_result)

    # Direct URL Fetching: if message contains URLs, fetch them directly
    urls_in_msg = re.findall(r"https?://[^\s]+", message)
    fetched_urls = set()
    for u in urls_in_msg:
        u_clean = u.rstrip(".,;!?()\"'")
        if u_clean in fetched_urls:
            continue
        fetched_urls.add(u_clean)
        try:
            print(f"[DIRECT URL FETCH] Fetching page content for: {u_clean}")
            async with httpx.AsyncClient(timeout=5.0, headers={"User-Agent": USER_AGENT}) as client:
                resp = await client.get(u_clean)
                resp.raise_for_status()
                body = _extract_main_body(resp.text)
                if body:
                    context_parts.append(
                        sanitize_untrusted_text(
                            f"--- Content of URL: {u_clean} ---\n{body[:1500]}\n"
                        )
                    )
        except Exception as e:
            print(f"⚠️ Direct URL fetch failed for {u_clean}: {e}")

    search_query = _extract_search_query(message)
    if not search_query and _should_search_web(message):
        # Only search web if there are no URLs in the message, to prevent search engine noise
        if not urls_in_msg:
            search_query = message.strip()

    if search_query:
        try:
            results = await search_web(search_query)
            # Check for missing player transfer story
            is_missing, team, rival_team = await detect_missing_player_in_transfer(search_query, results)
            if is_missing and team and rival_team:
                followup_query = f"{team} transfer {rival_team} player 2025"
                print(f"[FOLLOW-UP SEARCH] Running second search: {followup_query}")
                try:
                    followup_results = await search_web(followup_query)
                    has_player, player_name = await check_if_player_found(followup_query, followup_results)
                    if has_player and player_name:
                        context_parts.append(format_search_results(followup_results, followup_query))
                    else:
                        source_url = results[0]["url"] if results else "https://lite.duckduckgo.com"
                        context_parts.append(f"I found the story but couldn't confirm the player name. Here is the source: {source_url}")
                except Exception as followup_exc:
                    print(f"⚠️ Follow-up search failed: {followup_exc}")
                    source_url = results[0]["url"] if results else "https://lite.duckduckgo.com"
                    context_parts.append(f"I found the story but couldn't confirm the player name. Here is the source: {source_url}")
            else:
                context_parts.append(format_search_results(results, search_query))
        except Exception as exc:
            context_parts.append(sanitize_untrusted_text(f"Web search failed for '{search_query}': {exc}"))

    artifact_kind, artifact_title = _extract_artifact_request(message)
    if artifact_kind:
        artifact_message = sanitize_untrusted_text(
            (
                f"Document generation request detected: {artifact_kind.upper()}.\n"
                "Create the requested content directly. Do not claim you cannot generate files."
            )
        )
        context_parts.append(artifact_message)

    return ToolBundle(
        context="\n".join(context_parts).strip(),
        artifact_kind=artifact_kind,
        artifact_title=artifact_title,
    )
